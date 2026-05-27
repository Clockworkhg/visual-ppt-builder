#!/usr/bin/env python3
"""Build an editable PPTX from a visual-ppt-builder slide_plan.json.

This helper covers two paths:
- background_plus_text: optional per-slide background images plus text boxes
- reconstruct_editable: native shapes, independent image assets, and text boxes

For high-polish decks, use the Presentations skill/plugin when available.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - environment dependent
    raise SystemExit(
        "python-pptx and Pillow are required for build_ppt.py. "
        "Install them or use the Presentations skill/plugin."
    ) from exc


SLIDE_W = 13.333
SLIDE_H = 7.5

SHAPE_KINDS = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
}


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def hex_to_rgb(value: str | None, fallback: str = "#222222") -> RGBColor:
    raw = (value or fallback).strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def palette(plan: dict[str, Any]) -> dict[str, str]:
    style = plan.get("style") or {}
    pal = style.get("palette") or {}
    return {
        "background": pal.get("background", "#F7F8FA"),
        "primary": pal.get("primary", "#2F7DF6"),
        "secondary": pal.get("secondary", "#F2C94C"),
        "text": pal.get("text", "#182026"),
        "muted": pal.get("muted", "#667085"),
        "panel": pal.get("panel", "#FFFFFF"),
    }


def font_name(plan: dict[str, Any]) -> str:
    style = plan.get("style") or {}
    return style.get("font") or "Microsoft YaHei"


def asset_lookup(plan: dict[str, Any]) -> dict[str, dict[str, Any]]:
    result: dict[str, dict[str, Any]] = {}
    for item in plan.get("assets") or []:
        if isinstance(item, dict) and item.get("id"):
            result[str(item["id"])] = item
    return result


def resolve_path(path_value: str, plan_dir: Path, images_dir: Path | None) -> Path | None:
    raw = Path(path_value)
    candidates: list[Path] = []
    if raw.is_absolute():
        candidates.append(raw)
    else:
        candidates.append(plan_dir / raw)
        if images_dir:
            candidates.append(images_dir / raw.name)
            candidates.append(images_dir / raw)
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def find_background(
    slide_spec: dict[str, Any],
    plan_dir: Path,
    images_dir: Path | None,
) -> Path | None:
    candidates: list[Path] = []
    bg = slide_spec.get("background_image")
    if bg:
        resolved = resolve_path(str(bg), plan_dir, images_dir)
        if resolved:
            return resolved

    index = int(slide_spec.get("index", 0) or 0)
    if images_dir and index:
        for ext in (".png", ".jpg", ".jpeg", ".webp"):
            candidates.append(images_dir / f"slide_{index:02d}{ext}")
            candidates.append(images_dir / f"slide-{index:02d}{ext}")

    for path in candidates:
        if path.exists():
            return path
    return None


def add_native_background(slide: Any, pal: dict[str, str], slide_type: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(pal["background"], "#F7F8FA")

    accent = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.6),
        Inches(-0.7),
        Inches(4.8),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(pal["primary"], "#2F7DF6")
    accent.fill.transparency = 72
    accent.line.fill.background()

    if slide_type not in {"cover", "closing", "section"}:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.16),
            Inches(SLIDE_H),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(pal["primary"], "#2F7DF6")
        bar.line.fill.background()


def add_plan_shape(slide: Any, spec: dict[str, Any], pal: dict[str, str]) -> None:
    kind = str(spec.get("kind", "rect")).lower()
    if kind == "line":
        start = spec.get("start") or {}
        end = spec.get("end") or {}
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(float(start.get("x", 0))),
            Inches(float(start.get("y", 0))),
            Inches(float(end.get("x", 1))),
            Inches(float(end.get("y", 1))),
        )
        line = spec.get("line", pal["primary"])
        connector.line.color.rgb = hex_to_rgb(str(line), pal["primary"])
        connector.line.width = Pt(float(spec.get("line_width", 1.5)))
        return

    shape_type = SHAPE_KINDS.get(kind)
    if not shape_type:
        raise ValueError(f"Unsupported native shape kind: {kind}")

    box = spec.get("box") or {}
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(float(box.get("x", 0))),
        Inches(float(box.get("y", 0))),
        Inches(float(box.get("w", 1))),
        Inches(float(box.get("h", 1))),
    )

    fill = spec.get("fill", pal["primary"])
    if str(fill).lower() in {"none", "transparent"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(str(fill), pal["primary"])
        shape.fill.transparency = int(float(spec.get("transparency", 0)))

    line = spec.get("line", "none")
    if line is None or str(line).lower() in {"none", "transparent"}:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(str(line), pal["muted"])
        shape.line.width = Pt(float(spec.get("line_width", 1)))

    if "rotation" in spec:
        shape.rotation = float(spec["rotation"])


def add_image_asset(
    slide: Any,
    spec: dict[str, Any],
    assets: dict[str, dict[str, Any]],
    plan_dir: Path,
    images_dir: Path | None,
    warnings: list[str],
    slide_pos: int,
) -> None:
    path_value = spec.get("path")
    asset_id = spec.get("asset_id")
    if not path_value and asset_id:
        asset = assets.get(str(asset_id))
        if asset:
            path_value = asset.get("path")

    if not path_value:
        warnings.append(f"Slide {slide_pos}: image asset missing path or asset_id.")
        return

    image_path = resolve_path(str(path_value), plan_dir, images_dir)
    if not image_path:
        warnings.append(f"Slide {slide_pos}: image asset not found: {path_value}")
        return

    box = spec.get("box") or {}
    fit = str(spec.get("fit", "contain")).lower()
    if fit not in {"contain", "cover", "stretch"}:
        warnings.append(
            f"Slide {slide_pos}: unsupported image fit {fit!r}; using contain."
        )
        fit = "contain"

    x = float(box.get("x", 0))
    y = float(box.get("y", 0))
    w = float(box.get("w", 1))
    h = float(box.get("h", 1))

    with Image.open(image_path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    box_ratio = w / h if h else image_ratio

    if fit == "stretch":
        picture = slide.shapes.add_picture(
            str(image_path),
            Inches(x),
            Inches(y),
            width=Inches(w),
            height=Inches(h),
        )
    elif fit == "cover":
        picture = slide.shapes.add_picture(
            str(image_path),
            Inches(x),
            Inches(y),
            width=Inches(w),
            height=Inches(h),
        )
        if image_ratio > box_ratio:
            crop = max(0, min((1 - box_ratio / image_ratio) / 2, 0.499))
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < box_ratio:
            crop = max(0, min((1 - image_ratio / box_ratio) / 2, 0.499))
            picture.crop_top = crop
            picture.crop_bottom = crop
    else:
        if image_ratio > box_ratio:
            draw_w = w
            draw_h = w / image_ratio
            draw_x = x
            draw_y = y + (h - draw_h) / 2
        else:
            draw_h = h
            draw_w = h * image_ratio
            draw_x = x + (w - draw_w) / 2
            draw_y = y
        picture = slide.shapes.add_picture(
            str(image_path),
            Inches(draw_x),
            Inches(draw_y),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    if "rotation" in spec:
        picture.rotation = float(spec["rotation"])


def add_readability_panel(slide: Any, pal: dict[str, str], slide_type: str) -> None:
    if slide_type in {"section"}:
        return
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(0.55 if slide_type in {"cover", "closing"} else 0.35),
        Inches(6.2 if slide_type in {"cover", "closing"} else 6.6),
        Inches(5.6 if slide_type in {"cover", "closing"} else 6.45),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = hex_to_rgb(pal["panel"], "#FFFFFF")
    panel.fill.transparency = 12
    panel.line.fill.background()


def clear_text_frame(text_frame: Any) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.03)
    text_frame.margin_right = Inches(0.03)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)


def add_text_box(
    slide: Any,
    text: str,
    box: dict[str, float],
    size: float,
    color: str,
    family: str,
    bold: bool = False,
    align: str = "left",
) -> None:
    shape = slide.shapes.add_textbox(
        Inches(float(box["x"])),
        Inches(float(box["y"])),
        Inches(float(box["w"])),
        Inches(float(box["h"])),
    )
    tf = shape.text_frame
    clear_text_frame(tf)

    lines = str(text).splitlines() or [""]
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = {
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = family
        run.font.size = Pt(float(size))
        run.font.bold = bool(bold)
        run.font.color.rgb = hex_to_rgb(color, "#222222")


def default_blocks(slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    slide_type = str(slide_spec.get("type", "content")).lower()
    title = slide_spec.get("title", "")
    subtitle = slide_spec.get("subtitle") or slide_spec.get("tagline")
    body = slide_spec.get("body") or slide_spec.get("bullets")
    blocks: list[dict[str, Any]] = []

    if slide_type == "section":
        blocks.append(
            {
                "role": "title",
                "text": title,
                "box": {"x": 1.1, "y": 2.8, "w": 11.1, "h": 1.1},
                "font_size": 36,
                "align": "center",
            }
        )
        if subtitle:
            blocks.append(
                {
                    "role": "body",
                    "text": subtitle,
                    "box": {"x": 2.0, "y": 4.0, "w": 9.3, "h": 0.8},
                    "font_size": 17,
                    "align": "center",
                }
            )
        return blocks

    if slide_type in {"cover", "closing"}:
        blocks.append(
            {
                "role": "title",
                "text": title,
                "box": {"x": 0.75, "y": 1.25, "w": 5.95, "h": 1.25},
                "font_size": 38,
            }
        )
        if subtitle:
            blocks.append(
                {
                    "role": "subtitle",
                    "text": subtitle,
                    "box": {"x": 0.82, "y": 2.65, "w": 5.55, "h": 0.8},
                    "font_size": 18,
                }
            )
        return blocks

    blocks.append(
        {
            "role": "title",
            "text": title,
            "box": {"x": 0.72, "y": 0.52, "w": 6.1, "h": 0.75},
            "font_size": 26,
        }
    )
    if body:
        if isinstance(body, list):
            text = "\n".join(f"- {item}" for item in body)
        else:
            text = str(body)
        blocks.append(
            {
                "role": "body",
                "text": text,
                "box": {"x": 0.78, "y": 1.55, "w": 5.8, "h": 4.7},
                "font_size": 16,
            }
        )
    return blocks


def normalize_blocks(slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    raw_blocks = slide_spec.get("text_blocks")
    if not raw_blocks:
        return default_blocks(slide_spec)

    blocks: list[dict[str, Any]] = []
    for idx, block in enumerate(raw_blocks):
        if isinstance(block, str):
            role = "title" if idx == 0 else "body"
            blocks.append({"role": role, "text": block})
        elif isinstance(block, dict):
            blocks.append(dict(block))

    if not any((b.get("role") == "title") for b in blocks) and slide_spec.get("title"):
        blocks.insert(0, {"role": "title", "text": slide_spec["title"]})
    return blocks


def default_box(role: str, slide_type: str) -> dict[str, float]:
    if slide_type in {"cover", "closing"}:
        if role == "title":
            return {"x": 0.75, "y": 1.25, "w": 5.95, "h": 1.25}
        if role == "subtitle":
            return {"x": 0.82, "y": 2.65, "w": 5.55, "h": 0.8}
        return {"x": 0.82, "y": 3.55, "w": 5.55, "h": 2.0}
    if slide_type == "section":
        return {"x": 1.1, "y": 2.8, "w": 11.1, "h": 1.1}
    if role == "title":
        return {"x": 0.72, "y": 0.52, "w": 6.1, "h": 0.75}
    return {"x": 0.78, "y": 1.55, "w": 5.8, "h": 4.7}


def estimate_overflow(text: str, box: dict[str, float], size: float) -> bool:
    width_pt = float(box["w"]) * 72
    height_pt = float(box["h"]) * 72
    chars_per_line = max(int(width_pt / max(size * 0.55, 1)), 1)
    line_capacity = max(int(height_pt / max(size * 1.35, 1)), 1)
    hard_lines = max(len(str(text).splitlines()), 1)
    estimated_lines = 0
    for line in str(text).splitlines() or [""]:
        estimated_lines += max((len(line) + chars_per_line - 1) // chars_per_line, 1)
    return max(estimated_lines, hard_lines) > line_capacity


def build_deck(plan: dict[str, Any], plan_path: Path, images_dir: Path | None) -> tuple[Presentation, list[str]]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    pal = palette(plan)
    family = font_name(plan)
    assets = asset_lookup(plan)
    warnings: list[str] = []

    slides = plan.get("slides") or []
    for pos, slide_spec in enumerate(slides, start=1):
        slide_type = str(slide_spec.get("type", "content")).lower()
        slide = prs.slides.add_slide(blank)
        bg = find_background(slide_spec, plan_path.parent, images_dir)
        if bg:
            slide.shapes.add_picture(str(bg), Inches(0), Inches(0), width=Inches(SLIDE_W), height=Inches(SLIDE_H))
            if plan.get("mode") != "reconstruct_editable":
                add_readability_panel(slide, pal, slide_type)
        else:
            add_native_background(slide, pal, slide_type)

        for shape_spec in slide_spec.get("native_shapes") or []:
            if isinstance(shape_spec, dict):
                add_plan_shape(slide, shape_spec, pal)

        for image_spec in slide_spec.get("image_assets") or []:
            if isinstance(image_spec, dict):
                add_image_asset(
                    slide,
                    image_spec,
                    assets,
                    plan_path.parent,
                    images_dir,
                    warnings,
                    pos,
                )

        for block in normalize_blocks(slide_spec):
            role = str(block.get("role", "body")).lower()
            text = str(block.get("text", ""))
            if not text.strip():
                continue
            box = block.get("box") or default_box(role, slide_type)
            size = float(block.get("font_size") or (34 if role == "title" else 16))
            color = block.get("color") or (pal["text"] if role == "title" else pal["muted"])
            bold = bool(block.get("bold", role == "title"))
            align = str(block.get("align", "left")).lower()
            add_text_box(slide, text, box, size, color, family, bold=bold, align=align)
            if estimate_overflow(text, box, size):
                warnings.append(
                    f"Slide {pos}: possible text overflow in {role!r} block."
                )

    return prs, warnings


def write_report(
    path: Path,
    plan: dict[str, Any],
    out_path: Path,
    warnings: list[str],
    images_dir: Path | None,
) -> None:
    slide_count = len(plan.get("slides") or [])
    lines = [
        "# Build Report",
        "",
        f"- Title: {plan.get('title', 'Untitled')}",
        f"- Mode: {plan.get('mode', 'editable')}",
        f"- Slide count: {slide_count}",
        f"- Output: {out_path}",
        f"- Images directory: {images_dir if images_dir else 'not provided'}",
        f"- Warnings: {len(warnings)}",
        f"- Top-level assets: {len(plan.get('assets') or [])}",
        f"- Native shape placements: {sum(len(s.get('native_shapes') or []) for s in plan.get('slides') or [])}",
        f"- Image asset placements: {sum(len(s.get('image_assets') or []) for s in plan.get('slides') or [])}",
        "",
    ]
    if warnings:
        lines.append("## Warnings")
        lines.extend(f"- {item}" for item in warnings)
        lines.append("")
    lines.extend(
        [
            "## Editability",
            "",
            "The helper inserts slide copy as PowerPoint text boxes. Generated",
            "or supplied visual assets are inserted below the text layer.",
            "",
        ]
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--plan", required=True, type=Path, help="Path to slide_plan.json")
    parser.add_argument("--images", type=Path, help="Directory containing slide_01.png, etc.")
    parser.add_argument("--out", required=True, type=Path, help="Output PPTX path")
    parser.add_argument("--report", type=Path, help="Build report markdown path")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    plan = load_json(args.plan)
    prs, warnings = build_deck(plan, args.plan, args.images)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    if args.report:
        write_report(args.report, plan, args.out, warnings, args.images)
    print(f"Built {args.out} with {len(plan.get('slides') or [])} slides.")
    if warnings:
        print(f"Warnings: {len(warnings)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
